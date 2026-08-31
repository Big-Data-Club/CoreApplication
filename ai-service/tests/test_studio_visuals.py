from __future__ import annotations

import asyncio
from io import BytesIO

from app.agents.tools.teacher.mcp_generate_slide_deck import McpGenerateSlideDeckTool
from app.agents.tools.teacher.mcp_generate_report import McpGenerateReportTool
from app.services.studio.doc_renderer import render_plan_to_markdown
from app.services.studio.plan_schema import PlanSection, StudioPlan
from app.services.studio.plan_schema import coerce_plan
from app.services.studio.pptx_renderer import render_plan_to_pptx_bytes
from app.services.studio.visuals import mermaid_for_visual


def _plan() -> StudioPlan:
    return StudioPlan(
        title="Visual lesson",
        sections=[PlanSection(
            title="Request lifecycle",
            slide_bullets=["Client sends request", "Server validates", "LMS stores draft"],
            visual_type="flow",
            visual_labels=["Client", "Validation", "Draft"],
        ).cleaned()],
    )


def test_mermaid_is_bounded_and_sanitized():
    diagram = mermaid_for_visual("flow", ['A[bad]", click X', "Safe result"])
    assert "click" in diagram  # treated as inert label text, never a Mermaid directive
    assert "[bad]" not in diagram
    assert diagram.startswith("flowchart LR")


def test_studio_markdown_contains_visual():
    markdown = render_plan_to_markdown(_plan())
    assert "```mermaid" in markdown
    assert "Client" in markdown


def test_invalid_model_visual_type_falls_back_safely():
    plan, warnings = coerce_plan(
        {"title": "x", "sections": [{"title": "y", "slide_bullets": ["A", "B"], "visual_type": "javascript"}]},
        kind="slides", fallback_title="", target_sections=1,
    )
    assert warnings == []
    assert plan is not None
    assert plan.sections[0].visual_type == "auto"
    assert plan.sections[0].visual_labels == ["A", "B"]


def test_studio_pptx_contains_editable_visual_shapes():
    from pptx import Presentation
    presentation = Presentation(BytesIO(render_plan_to_pptx_bytes(_plan())))
    section_slide = presentation.slides[1]
    texts = [shape.text for shape in section_slide.shapes if hasattr(shape, "text_frame")]
    assert any("Client" in text for text in texts)
    assert any("Validation" in text for text in texts)


def test_mcp_deck_has_visual_on_every_slide():
    result = asyncio.run(McpGenerateSlideDeckTool().execute(
        title="MCP visuals",
        slides=[{"title": "Flow", "bullets": ["One", "Two"], "visual_type": "flow", "visual_labels": ["One", "Two"]}],
    ))
    assert result.status == "success"
    assert result.data["visual_coverage"] == "1/1"
    assert "```mermaid" in result.data["reveal_markdown"]


def test_mcp_report_has_grounding_visual_and_non_blocking_image_prompt():
    result = asyncio.run(McpGenerateReportTool().execute(
        title="Architecture report",
        executive_summary="A grounded summary.",
        sections=[{
            "heading": "Data flow",
            "body_markdown": "The event crosses the validated pipeline.",
            "key_findings": ["Validation precedes persistence"],
            "visual_type": "flow",
            "visual_labels": ["Client", "Validation", "Storage"],
            "alt_text": "Three stages from client to storage.",
            "illustration_prompt": "Minimal academic pipeline illustration",
            "source_refs": ["LOCAL-01 §Data flow", "LMS-S3"],
        }],
    ))
    assert result.status == "success"
    assert result.data["visual_coverage"] == "1/1"
    markdown = result.data["report_markdown"]
    assert "```mermaid" in markdown
    assert "Gợi ý ảnh (chưa tạo)" in markdown
    assert "LOCAL-01" in markdown
