"""
ai-service/app/services/studio/pptx_renderer.py

Native-PPTX slide renderer (python-pptx). Slides are REAL editable objects -
never baked images - so teachers can fix a typo in PowerPoint afterwards.

Themes: 'academic' (navy/white), 'modern' (dark cyan), 'minimal' (light).
Watermark: HPCC logo bottom-left, BDC logo bottom-right on every slide
(assets copied from the frontend brand kit).
"""
from __future__ import annotations

import logging
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

from app.services.studio.plan_schema import StudioPlan
from app.services.studio.visuals import resolve_visual_type

logger = logging.getLogger(__name__)

_BRAND_DIR = Path(__file__).resolve().parents[2] / "assets" / "brand"
_HPCC_LOGO = _BRAND_DIR / "hpcc-logo.png"
_BDC_LOGO = _BRAND_DIR / "bdclogo_backup.png"

SLIDE_W = Inches(13.333)   # 16:9
SLIDE_H = Inches(7.5)

THEMES: dict[str, dict] = {
    "academic": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title": RGBColor(0x1B, 0x2C, 0x4E),
        "accent": RGBColor(0x25, 0x54, 0xD2),
        "text": RGBColor(0x33, 0x41, 0x55),
        "muted": RGBColor(0x8A, 0x94, 0xA6),
    },
    "modern": {
        "bg": RGBColor(0x0F, 0x1E, 0x35),
        "title": RGBColor(0xFF, 0xFF, 0xFF),
        "accent": RGBColor(0x22, 0xD3, 0xEE),
        "text": RGBColor(0xDB, 0xE4, 0xF0),
        "muted": RGBColor(0x7C, 0x8B, 0xA5),
    },
    "minimal": {
        "bg": RGBColor(0xFA, 0xFA, 0xF7),
        "title": RGBColor(0x22, 0x22, 0x22),
        "accent": RGBColor(0xE2, 0x62, 0x2B),
        "text": RGBColor(0x3D, 0x3D, 0x3D),
        "muted": RGBColor(0x9E, 0x9E, 0x9E),
    },
}


def _set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(slide, left, top, width, height, text, *, size, color,
              bold=False, align=PP_ALIGN.LEFT, font="Segoe UI"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return tf


def _add_bullets(tf, bullets: list[str], *, size, color, accent):
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 and not tf.paragraphs[0].runs else tf.add_paragraph()
        p.space_after = Pt(10)
        dash = p.add_run()
        dash.text = "▸ "
        dash.font.color.rgb = accent
        dash.font.size = Pt(size)
        run = p.add_run()
        run.text = b
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Segoe UI"


def _watermark(slide) -> None:
    for path, (w_in, left_in, top_in) in (
        (_HPCC_LOGO, (1.05, SLIDE_W.inches - 1.15, SLIDE_H.inches - 0.42)),
        (_BDC_LOGO, (0.55, SLIDE_W.inches - 0.62, SLIDE_H.inches - 0.44)),
    ):
        try:
            if path.exists():
                slide.shapes.add_picture(
                    str(path), Emu(int(left_in * 914400)), Emu(int(top_in * 914400)),
                    width=Emu(int(w_in * 914400)),
                )
        except Exception as exc:  # noqa: BLE001 - watermark must never break render
            logger.warning("watermark skip %s: %s", path.name, exc)


def _style_visual_shape(shape, text: str, *, fill: RGBColor, text_color: RGBColor,
                        font_size: int = 14) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Segoe UI"
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = text_color


def _add_native_visual(slide, sec, theme: dict) -> None:
    """Draw an editable infographic; no network fetch or image service required."""
    labels = list(sec.visual_labels or [])[:6]
    if not labels:
        return
    visual_type = resolve_visual_type(sec.visual_type, labels)
    left, top, width, height = 7.55, 1.65, 5.05, 4.75
    panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    panel.fill.background()
    panel.line.color.rgb = theme["accent"]
    panel.line.width = Pt(1.5)

    white = RGBColor(0xFF, 0xFF, 0xFF)
    if visual_type == "comparison":
        cols = 2
        card_w = 2.05
        card_h = 1.15
        for idx, label in enumerate(labels[:4]):
            x = left + 0.35 + (idx % cols) * 2.3
            y = top + 0.55 + (idx // cols) * 1.65
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_w), Inches(card_h))
            _style_visual_shape(shape, label, fill=theme["accent"], text_color=white)
    elif visual_type == "hierarchy":
        root = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left + 1.35), Inches(top + 0.4), Inches(2.35), Inches(0.9))
        _style_visual_shape(root, labels[0], fill=theme["accent"], text_color=white)
        children = labels[1:5] or labels[:1]
        child_w = min(1.95, 4.3 / max(1, len(children)))
        for idx, label in enumerate(children):
            x = left + 0.3 + idx * (4.45 / max(1, len(children)))
            child = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(top + 2.25), Inches(child_w), Inches(1.25))
            _style_visual_shape(child, label, fill=theme["accent"], text_color=white, font_size=12)
    elif visual_type == "cycle":
        positions = [(1.7, 0.35), (3.25, 1.55), (2.55, 3.0), (0.75, 3.0), (0.05, 1.55)]
        for idx, label in enumerate(labels[:5]):
            x, y = positions[idx]
            shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left + x), Inches(top + y), Inches(1.45), Inches(0.95))
            _style_visual_shape(shape, label, fill=theme["accent"], text_color=white, font_size=11)
    else:
        horizontal = visual_type in {"flow", "timeline"} and len(labels) <= 4
        if horizontal:
            card_w = 3.95 / len(labels)
            for idx, label in enumerate(labels):
                x = left + 0.25 + idx * (4.55 / len(labels))
                shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(top + 1.45), Inches(card_w), Inches(1.25))
                _style_visual_shape(shape, label, fill=theme["accent"], text_color=white, font_size=11)
                if idx < len(labels) - 1:
                    arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x + card_w + 0.05), Inches(top + 1.82), Inches(0.28), Inches(0.45))
                    arrow.fill.solid(); arrow.fill.fore_color.rgb = theme["muted"]; arrow.line.fill.background()
        else:
            card_h = min(0.68, 3.45 / len(labels))
            for idx, label in enumerate(labels):
                y = top + 0.38 + idx * (3.85 / len(labels))
                shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left + 0.55), Inches(y), Inches(3.95), Inches(card_h))
                _style_visual_shape(shape, label, fill=theme["accent"], text_color=white, font_size=11)

    if sec.visual_suggestion:
        _add_text(slide, Inches(left + 0.3), Inches(top + height - 0.55), Inches(width - 0.6), Inches(0.35),
                  sec.visual_suggestion, size=9, color=theme["muted"], align=PP_ALIGN.CENTER)


def render_plan_to_pptx_bytes(plan: StudioPlan, theme: str = "academic") -> bytes:
    t = THEMES.get(theme, THEMES["academic"])
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    def new_slide():
        s = prs.slides.add_slide(blank)
        _set_bg(s, t["bg"])
        return s

    # ── Title slide ────────────────────────────────────────────────────────
    title_slide = new_slide()
    bar = title_slide.shapes.add_shape(1, Inches(0), Inches(3.05), SLIDE_W, Pt(6))
    bar.fill.solid(); bar.fill.fore_color.rgb = t["accent"]; bar.line.fill.background()
    _add_text(title_slide, Inches(0.9), Inches(2.0), SLIDE_W - Inches(1.8), Inches(1.0),
              plan.safe_title, size=40, color=t["title"], bold=True)
    if plan.learning_objectives:
        _add_text(title_slide, Inches(0.9), Inches(3.35), SLIDE_W - Inches(1.8), Inches(2.4),
                  "Mục tiêu: " + "; ".join(plan.learning_objectives[:4]),
                  size=16, color=t["muted"])
    _watermark(title_slide)

    # ── Agenda (optional, when many sections) ─────────────────────────────
    if len(plan.sections) >= 4:
        agenda = new_slide()
        _add_text(agenda, Inches(0.9), Inches(0.55), Inches(11), Inches(0.8),
                  "Nội dung", size=32, color=t["title"], bold=True)
        tf = _add_text(agenda, Inches(1.0), Inches(1.6), Inches(11.3), Inches(5.2),
                       "", size=18, color=t["text"])
        _add_bullets(tf, [f"{i+1}. {s.title}" for i, s in enumerate(plan.sections)],
                     size=18, color=t["text"], accent=t["accent"])
        _watermark(agenda)

    # ── Section slides ─────────────────────────────────────────────────────
    for idx, sec in enumerate(plan.sections, start=1):
        slide = new_slide()

        # Accent rail
        rail = slide.shapes.add_shape(1, Inches(0), Inches(0), Pt(8), SLIDE_H)
        rail.fill.solid(); rail.fill.fore_color.rgb = t["accent"]; rail.line.fill.background()

        _add_text(slide, Inches(0.75), Inches(0.45), Inches(11.6), Inches(0.9),
                  f"{idx}. {sec.title}", size=30, color=t["title"], bold=True)

        bullets = sec.slide_bullets or sec.key_points or []
        body_tf = _add_text(slide, Inches(1.0), Inches(1.7), Inches(5.95), Inches(4.6),
                            "", size=20, color=t["text"])
        if bullets:
            _add_bullets(body_tf, bullets, size=20, color=t["text"], accent=t["accent"])
        elif sec.narration:
            head = sec.narration[:600]
            _add_bullets(body_tf, [head], size=16, color=t["text"], accent=t["accent"])

        _add_native_visual(slide, sec, t)

        # Slide number chip
        chip = slide.shapes.add_shape(1, Inches(12.55), Inches(0.35), Inches(0.6), Inches(0.35))
        chip.fill.solid(); chip.fill.fore_color.rgb = t["accent"]; chip.line.fill.background()
        ctf = chip.text_frame; ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
        crun = ctf.paragraphs[0].add_run(); crun.text = str(idx)
        crun.font.size = Pt(12); crun.font.bold = True
        crun.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        _watermark(slide)

    # ── Closing slide ──────────────────────────────────────────────────────
    end = new_slide()
    _add_text(end, Inches(0.9), Inches(2.9), SLIDE_W - Inches(1.8), Inches(1.2),
              "Cảm ơn các bạn đã theo dõi!", size=36, color=t["title"], bold=True,
              align=PP_ALIGN.CENTER)
    if plan.summary:
        _add_text(end, Inches(1.6), Inches(4.2), SLIDE_W - Inches(3.2), Inches(1.6),
                  plan.summary, size=14, color=t["muted"], align=PP_ALIGN.CENTER)
    _watermark(end)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def render_section_to_pptx_bytes(section, *, deck_title: str, index: int,
                                 total: int, theme: str = "academic") -> bytes:
    """Single-section mini-deck used for per-section re-render previews."""
    single = StudioPlan(
        title=f"{deck_title} — {section.title}", sections=[section],
        language="vi", learning_objectives=[], summary="",
    )
    data = render_plan_to_pptx_bytes(single, theme=theme)
    return data
