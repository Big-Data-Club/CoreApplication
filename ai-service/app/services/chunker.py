from __future__ import annotations

import io
import logging
import re
from dataclasses import dataclass

logger = logging.getLogger(__name__)


@dataclass
class DocumentChunk:
    text: str
    index: int
    source_type: str           # 'document' | 'video'
    page_number: int | None = None
    start_time_sec: int | None = None
    end_time_sec: int | None = None
    language: str = "vi"


def sanitize_text(text: str) -> str:
    """
    Remove characters PostgreSQL UTF-8 cannot store.
    - Null bytes (0x00): extracted by PyMuPDF from some PDFs
    - Other non-printable control chars (keep \\n \\r \\t which are fine)
    """
    return re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)


class PDFChunker:
    """
    Chunk PDF files page-by-page, then split large pages into smaller chunks.
    Preserves page_number for Deep Link.
    """

    def __init__(self, chunk_size: int = 500, overlap: int = 50):
        self.chunk_size = chunk_size
        self.overlap = overlap

    def chunk_bytes(self, pdf_bytes: bytes) -> list[DocumentChunk]:
        """Process raw PDF bytes → list of chunks with page metadata."""
        try:
            return self._chunk_with_pymupdf4llm(pdf_bytes)
        except Exception as e:
            logger.warning(f"pymupdf4llm failed ({e}), falling back to PyMuPDF")
            try:
                return self._chunk_with_pymupdf(pdf_bytes)
            except Exception as e2:
                logger.warning(f"PyMuPDF failed ({e2}), falling back to pypdf")
                return self._chunk_with_pypdf(pdf_bytes)

    # Tier 1: pymupdf4llm 

    def _chunk_with_pymupdf4llm(self, pdf_bytes: bytes) -> list[DocumentChunk]:
        """
        Dùng pymupdf4llm để convert PDF → Markdown có cấu trúc.

        Ưu điểm so với text thuần:
        - Giữ headings (# ## ###) → chunk theo section
        - Bảng được render thành Markdown table
        - Bold/italic được giữ → LLM hiểu nhấn mạnh
        - Multi-column layout được xử lý đúng
        """
        import pymupdf4llm
        import pymupdf

        doc = pymupdf.open(stream=pdf_bytes, filetype="pdf")
        chunks: list[DocumentChunk] = []
        chunk_index = 0

        for page_num in range(len(doc)):
            # Extract markdown cho từng trang riêng lẻ để giữ page metadata
            try:
                md_text = pymupdf4llm.to_markdown(
                    doc,
                    pages=[page_num],
                    # Giữ heading structure
                    hdr_info=True,
                    # Bảng dạng markdown
                    table_strategy="lines_strict",
                )
            except TypeError:
                # Một số version cũ không hỗ trợ tất cả params
                md_text = pymupdf4llm.to_markdown(doc, pages=[page_num])

            md_text = sanitize_text(md_text.strip())
            if not md_text:
                continue

            # Chunk theo section (heading-aware)
            page_chunks = self._split_markdown_by_heading(md_text)
            for chunk_text in page_chunks:
                chunk_text = sanitize_text(chunk_text.strip())
                if len(chunk_text) < 20:  # bỏ chunk quá ngắn
                    continue
                chunks.append(DocumentChunk(
                    text=chunk_text,
                    index=chunk_index,
                    source_type="document",
                    page_number=page_num + 1,
                    language=detect_language(chunk_text),
                ))
                chunk_index += 1

        doc.close()
        return chunks

    # Tier 2: PyMuPDF + pdfplumber tables 

    def _chunk_with_pymupdf(self, pdf_bytes: bytes) -> list[DocumentChunk]:
        """
        PyMuPDF + pdfplumber cho bảng.
        Cho từng trang: extract text, detect bảng, merge.
        """
        import pymupdf

        chunks: list[DocumentChunk] = []
        chunk_index = 0
        doc = pymupdf.open(stream=pdf_bytes, filetype="pdf")

        for page_num in range(len(doc)):
            page = doc[page_num]
            text = sanitize_text(page.get_text("text").strip())

            # Thêm bảng từ pdfplumber nếu có
            tables_md = _try_extract_tables(pdf_bytes, page_num + 1)
            if tables_md:
                text = text + "\n\n" + "\n\n".join(tables_md)

            if not text.strip():
                continue

            page_chunks = self._split_text(text)
            for chunk_text in page_chunks:
                chunk_text = sanitize_text(chunk_text.strip())
                if chunk_text:
                    chunks.append(DocumentChunk(
                        text=chunk_text,
                        index=chunk_index,
                        source_type="document",
                        page_number=page_num + 1,
                        language=detect_language(chunk_text),
                    ))
                    chunk_index += 1

        doc.close()
        return chunks

    # Tier 3: pypdf 

    def _chunk_with_pypdf(self, pdf_bytes: bytes) -> list[DocumentChunk]:
        import pypdf

        chunks: list[DocumentChunk] = []
        chunk_index = 0
        reader = pypdf.PdfReader(io.BytesIO(pdf_bytes))

        for page_num, page in enumerate(reader.pages):
            text = sanitize_text((page.extract_text() or "").strip())
            if not text:
                continue

            page_chunks = self._split_text(text)
            for chunk_text in page_chunks:
                chunk_text = sanitize_text(chunk_text.strip())
                if chunk_text:
                    chunks.append(DocumentChunk(
                        text=chunk_text,
                        index=chunk_index,
                        source_type="document",
                        page_number=page_num + 1,
                        language=detect_language(chunk_text),
                    ))
                    chunk_index += 1

        return chunks

    # Chunking strategies 

    def _split_markdown_by_heading(self, md_text: str) -> list[str]:
        """
        Chia Markdown theo heading (# ## ###).
        Mỗi section là 1 chunk cơ bản, sau đó split nếu quá dài.

        Ví dụ:
          # Chương 1       → chunk riêng
          ## Định nghĩa    → chunk riêng (với context heading cha)
          ...text...       → phần của chunk heading gần nhất
        """
        # Tách theo heading pattern
        parts = re.split(r"(?m)^(#{1,3}\s+.+)$", md_text)
        if len(parts) <= 1:
            # Không có heading → split thông thường
            return self._split_text(md_text)

        chunks: list[str] = []
        current_h1 = ""
        current_h2 = ""
        current_text = ""

        for part in parts:
            part = part.strip()
            if not part:
                continue

            if re.match(r"^#\s+", part):
                # H1 heading
                if current_text:
                    chunks.extend(self._split_text(current_text))
                current_h1 = part.lstrip("# ").strip()
                current_h2 = ""
                current_text = part + "\n"

            elif re.match(r"^##\s+", part):
                # H2 heading
                if current_text:
                    chunks.extend(self._split_text(current_text))
                current_h2 = part.lstrip("# ").strip()
                # Thêm breadcrumb context: "Chương 1 > Định nghĩa"
                breadcrumb = ""
                if current_h1:
                    breadcrumb = f"{current_h1} > "
                current_text = breadcrumb + part + "\n"

            elif re.match(r"^###\s+", part):
                # H3 heading
                if current_text:
                    chunks.extend(self._split_text(current_text))
                breadcrumb = ""
                if current_h1 and current_h2:
                    breadcrumb = f"{current_h1} > {current_h2} > "
                elif current_h1:
                    breadcrumb = f"{current_h1} > "
                current_text = breadcrumb + part + "\n"

            else:
                # Nội dung text
                current_text = (current_text + "\n" + part).strip() + "\n"

                # Flush nếu quá dài
                if len(current_text) > self.chunk_size * 2:
                    # Giữ lại heading của section để chunk tiếp theo vẫn có context
                    heading_match = re.match(r"^(#{1,3}\s+.+)\n", current_text)
                    heading_prefix = heading_match.group(0) if heading_match else ""

                    split_chunks = self._split_text(current_text)
                    chunks.extend(split_chunks)
                    current_text = heading_prefix

        if current_text.strip():
            chunks.extend(self._split_text(current_text))

        return [c for c in chunks if len(c.strip()) >= 20]

    def _split_text(self, text: str) -> list[str]:
        """
        Split text into overlapping chunks at sentence boundaries.
        Prefers sentence breaks over hard character cuts.
        """
        if len(text) <= self.chunk_size:
            return [text] if text.strip() else []

        # Split theo sentence boundaries (VI + EN punctuation)
        sentences = re.split(r'(?<=[.!?。！？\n])\s+', text)
        chunks: list[str] = []
        current = ""

        for sentence in sentences:
            if len(current) + len(sentence) <= self.chunk_size:
                current = (current + " " + sentence).strip()
            else:
                if current:
                    chunks.append(current)
                if current and self.overlap > 0:
                    words = current.split()
                    overlap_text = " ".join(words[-max(1, self.overlap // 5):])
                    current = (overlap_text + " " + sentence).strip()
                else:
                    current = sentence

        if current:
            chunks.append(current)

        return chunks


class DocxChunker(PDFChunker):
    """
    Chunk Word (.docx) files.
    Cải tiến: chunk theo heading style (Heading 1/2/3) trước,
    giúp mỗi chunk tương ứng với một section có nghĩa.
    """

    def chunk_bytes(self, docx_bytes: bytes) -> list[DocumentChunk]:
        try:
            from docx import Document
        except ImportError:
            logger.error("python-docx not available")
            return []

        try:
            doc = Document(io.BytesIO(docx_bytes))
            sections: list[tuple[str, str]] = []  # (section_text, heading_context)
            current_heading = ""
            current_text_parts: list[str] = []

            for para in doc.paragraphs:
                if not para.text.strip():
                    continue

                # Phát hiện heading dựa trên style
                style_name = para.style.name if para.style else ""
                is_heading = (
                    "Heading" in style_name
                    or style_name.startswith("h")
                    or para.runs and para.runs[0].bold and len(para.text) < 100
                )

                if is_heading:
                    # Flush section hiện tại
                    if current_text_parts:
                        section_text = (current_heading + "\n" if current_heading else "") + "\n".join(current_text_parts)
                        sections.append((section_text, current_heading))
                    current_heading = para.text.strip()
                    current_text_parts = []
                else:
                    current_text_parts.append(para.text.strip())

            # Tables
            for table in doc.tables:
                rows = []
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    if any(row_text):
                        rows.append(row_text)
                if rows:
                    # Format bảng đơn giản
                    table_lines = [" | ".join(r) for r in rows]
                    current_text_parts.append("\n".join(table_lines))

            # Flush phần cuối
            if current_text_parts:
                section_text = (current_heading + "\n" if current_heading else "") + "\n".join(current_text_parts)
                sections.append((section_text, current_heading))

            # Build chunks
            chunks: list[DocumentChunk] = []
            chunk_index = 0
            for section_text, _ in sections:
                section_text = sanitize_text(section_text)
                if not section_text.strip():
                    continue
                for chunk_text in self._split_text(section_text):
                    chunk_text = sanitize_text(chunk_text.strip())
                    if len(chunk_text) >= 20:
                        chunks.append(DocumentChunk(
                            text=chunk_text,
                            index=chunk_index,
                            source_type="document",
                            page_number=1,
                            language=detect_language(chunk_text),
                        ))
                        chunk_index += 1

            return chunks

        except Exception as e:
            logger.error(f"Error chunking docx: {e}")
            return []


class PptxChunker(PDFChunker):
    """
    Chunk PowerPoint (.pptx) files.
    Mỗi slide là 1 unit tự nhiên → giữ nguyên, chỉ split nếu quá dài.
    Thêm title slide làm heading context cho các slide tiếp theo.
    """

    def chunk_bytes(self, pptx_bytes: bytes) -> list[DocumentChunk]:
        try:
            from pptx import Presentation
        except ImportError:
            logger.error("python-pptx not available")
            return []

        try:
            prs = Presentation(io.BytesIO(pptx_bytes))
            chunks: list[DocumentChunk] = []
            chunk_index = 0
            section_title = ""  # title của section/chapter slide

            for i, slide in enumerate(prs.slides):
                slide_parts: list[str] = []
                slide_title = ""

                # Lấy text từ shapes, ưu tiên title
                for shape in slide.shapes:
                    if not hasattr(shape, "text"):
                        continue
                    text = shape.text.strip()
                    if not text:
                        continue

                    # Phát hiện title shape
                    if (hasattr(shape, "placeholder_format")
                            and shape.placeholder_format
                            and shape.placeholder_format.idx == 0):
                        slide_title = text
                    else:
                        slide_parts.append(text)

                # Notes của slide
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes and len(notes) > 20:
                        slide_parts.append(f"[Ghi chú]: {notes}")

                if not slide_title and not slide_parts:
                    continue

                # Slide chỉ có title và ít text → có thể là section divider
                if slide_title and not slide_parts:
                    section_title = slide_title
                    # Vẫn tạo chunk để không mất thông tin
                    chunks.append(DocumentChunk(
                        text=slide_title,
                        index=chunk_index,
                        source_type="document",
                        page_number=i + 1,
                        language=detect_language(slide_title),
                    ))
                    chunk_index += 1
                    continue

                # Gộp toàn bộ text slide
                full_text = ""
                if section_title:
                    full_text += f"{section_title}\n"
                if slide_title:
                    full_text += f"{slide_title}\n"
                full_text += "\n".join(slide_parts)
                full_text = sanitize_text(full_text.strip())

                if not full_text:
                    continue

                for chunk_text in self._split_text(full_text):
                    chunk_text = sanitize_text(chunk_text.strip())
                    if len(chunk_text) >= 20:
                        chunks.append(DocumentChunk(
                            text=chunk_text,
                            index=chunk_index,
                            source_type="document",
                            page_number=i + 1,
                            language=detect_language(chunk_text),
                        ))
                        chunk_index += 1

            return chunks

        except Exception as e:
            logger.error(f"Error chunking pptx: {e}")
            return []


class ExcelChunker(PDFChunker):
    """
    Chunk Excel (.xlsx) files.
    Cải tiến: mỗi sheet là 1 context riêng, format bảng đẹp hơn.
    """

    def chunk_bytes(self, excel_bytes: bytes) -> list[DocumentChunk]:
        try:
            import openpyxl
        except ImportError:
            logger.error("openpyxl not available")
            return []

        try:
            wb = openpyxl.load_workbook(io.BytesIO(excel_bytes), data_only=True)
            all_chunks: list[DocumentChunk] = []
            chunk_index = 0

            for sheet in wb.worksheets:
                rows_data: list[list[str]] = []
                for row in sheet.iter_rows(values_only=True):
                    row_text = [str(cell) if cell is not None else "" for cell in row]
                    if any(t.strip() for t in row_text):
                        rows_data.append(row_text)

                if not rows_data:
                    continue

                # Format sheet thành text có cấu trúc
                sheet_header = f"Sheet: {sheet.title}"

                # Phát hiện header row (row đầu tiên thường là header)
                if len(rows_data) >= 2:
                    header_row = rows_data[0]
                    data_rows = rows_data[1:]
                    # Chia thành batch để chunk không quá lớn
                    batch_size = 20
                    for batch_start in range(0, len(data_rows), batch_size):
                        batch = data_rows[batch_start:batch_start + batch_size]
                        # Format theo kiểu "key: value" nếu nhiều cột
                        if len(header_row) > 1:
                            lines = [sheet_header]
                            for data_row in batch:
                                pairs = [
                                    f"{h}: {v}"
                                    for h, v in zip(header_row, data_row)
                                    if h.strip() and v.strip()
                                ]
                                if pairs:
                                    lines.append(" | ".join(pairs))
                        else:
                            lines = [sheet_header] + [" ".join(r) for r in batch if any(r)]

                        text = sanitize_text("\n".join(lines))
                        if text.strip():
                            all_chunks.append(DocumentChunk(
                                text=text,
                                index=chunk_index,
                                source_type="document",
                                page_number=1,
                                language=detect_language(text),
                            ))
                            chunk_index += 1
                else:
                    # Chỉ 1 row
                    text = sanitize_text(sheet_header + "\n" + " ".join(rows_data[0]))
                    if text.strip():
                        all_chunks.append(DocumentChunk(
                            text=text,
                            index=chunk_index,
                            source_type="document",
                            page_number=1,
                            language=detect_language(text),
                        ))
                        chunk_index += 1

            return all_chunks

        except Exception as e:
            logger.error(f"Error chunking excel: {e}")
            return []


class VideoTranscriptChunker:
    """
    Chunk video transcripts with timestamp metadata.
    Supports Whisper-style transcripts (SRT or JSON with timestamps).
    Groups segments into ~2-minute chunks for meaningful context.
    """

    def __init__(self, segment_duration_sec: int = 120, overlap_sec: int = 15):
        self.segment_duration = segment_duration_sec
        self.overlap = overlap_sec

    def chunk_whisper_json(self, transcript: dict) -> list[DocumentChunk]:
        segments = transcript.get("segments", [])
        if not segments:
            return []

        chunks: list[DocumentChunk] = []
        chunk_index = 0
        current_text = ""
        current_start = segments[0]["start"]
        current_end = segments[0]["end"]

        for seg in segments:
            segment_text = sanitize_text(seg.get("text", "").strip())
            seg_start = seg.get("start", 0)
            seg_end = seg.get("end", seg_start)

            if seg_end - current_start > self.segment_duration and current_text:
                chunks.append(DocumentChunk(
                    text=current_text.strip(),
                    index=chunk_index,
                    source_type="video",
                    start_time_sec=int(current_start),
                    end_time_sec=int(current_end),
                    language=detect_language(current_text),
                ))
                chunk_index += 1
                overlap_start = max(current_start, seg_end - self.overlap)
                current_text = segment_text
                current_start = overlap_start
                current_end = seg_end
            else:
                current_text = (current_text + " " + segment_text).strip()
                current_end = seg_end

        if current_text.strip():
            chunks.append(DocumentChunk(
                text=current_text.strip(),
                index=chunk_index,
                source_type="video",
                start_time_sec=int(current_start),
                end_time_sec=int(current_end),
                language=detect_language(current_text),
            ))

        return chunks

    def chunk_srt(self, srt_content: str) -> list[DocumentChunk]:
        blocks = re.split(r"\n\n+", srt_content.strip())
        segments = []
        for block in blocks:
            lines = block.strip().split("\n")
            if len(lines) < 2:
                continue
            for line in lines:
                match = re.match(
                    r"(\d{2}):(\d{2}):(\d{2}),(\d+)\s*-->\s*(\d{2}):(\d{2}):(\d{2}),(\d+)",
                    line,
                )
                if match:
                    h1, m1, s1, ms1, h2, m2, s2, ms2 = match.groups()
                    start = int(h1)*3600 + int(m1)*60 + int(s1)
                    end = int(h2)*3600 + int(m2)*60 + int(s2)
                    text_lines = [l for l in lines if not l.isdigit() and "-->" not in l]
                    text = sanitize_text(" ".join(text_lines).strip())
                    if text:
                        segments.append({"start": start, "end": end, "text": text})
                    break

        return self.chunk_whisper_json({"segments": segments})


def detect_language(text: str) -> str:
    """
    Simple language detection: if >5% of alphabetic characters are
    Vietnamese-specific diacritics, classify as Vietnamese.
    """
    vi_chars = set("àáảãạăắặằẳẵâấầẩẫậđèéẻẽẹêếềểễệìíỉĩịòóỏõọôốồổỗộơớờởỡợùúủũụưứừửữựỳýỷỹỵ"
                   "ÀÁẢÃẠĂẮẶẰẲẴÂẤẦẨẪẬĐÈÉẺẼẸÊẾỀỂỄỆÌÍỈĨỊÒÓỎÕỌÔỐỒỔỖỘƠỚỜỞỠỢÙÚỦŨỤƯỨỪỬỮỰỲÝỶỸỴ")
    total = len([c for c in text if c.isalpha()])
    if total == 0:
        return "vi"
    vi_count = sum(1 for c in text if c in vi_chars)
    return "vi" if vi_count / total > 0.05 else "en"


def format_timestamp(seconds: int) -> str:
    """Convert seconds to MM:SS or HH:MM:SS string."""
    h = seconds // 3600
    m = (seconds % 3600) // 60
    s = seconds % 60
    if h > 0:
        return f"{h:02d}:{m:02d}:{s:02d}"
    return f"{m:02d}:{s:02d}"

def _try_extract_tables(pdf_bytes: bytes, page_num: int) -> list[str]:
    """Trích xuất bảng bằng pdfplumber, không crash nếu không có."""
    try:
        from app.services.document_intelligence import extract_tables_from_page
        return extract_tables_from_page(pdf_bytes, page_num)
    except Exception:
        return []